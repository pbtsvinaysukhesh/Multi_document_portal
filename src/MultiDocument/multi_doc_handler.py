from docx import Document
from pypdf import PdfReader
import pandas as pd
from pptx import Presentation
import sqlite3
from sqlalchemy import create_engine
import logging
from pathlib import Path
from typing import Dict, Callable, Optional, List, Tuple
import sys
import os
import cv2
import pytesseract
from PIL import Image
import numpy as np
import camelot
import tabula
from pdf2image import convert_from_path
import tempfile

sys.path.append(os.path.join(os.path.dirname(__file__), '..', '..'))
from exception.custom_exception import (
    DocumentProcessingException,
    UnsupportedFormatException,
    FileNotFoundException,
    PDFProcessingException,
    DOCXProcessingException,
    ExcelProcessingException,
    CSVProcessingException,
    PowerPointProcessingException,
    TextProcessingException,
    MarkdownProcessingException,
    SQLiteProcessingException
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

class DocumentHandler:
    def __init__(self):
        self.supported_formats: Dict[str, Callable] = {
            '.docx': self.read_docx,
            '.pdf': self.read_pdf,
            '.txt': self.read_txt,
            '.xlsx': self.read_excel,
            '.csv': self.read_csv,
            '.ppt': self.read_ppt,
            '.pptx': self.read_ppt,
            '.md': self.read_md,
            '.db': self.read_sqlite
        }

    def read_file(self, file_path: str) -> Optional[str]:
        """
        Read and process a file based on its extension.

        Args:
            file_path: Path to the file to be processed

        Returns:
            Extracted text content from the file

        Raises:
            UnsupportedFormatException: If file format is not supported
            FileNotFoundException: If file doesn't exist
        """
        try:
            path = Path(file_path)
            if not path.exists():
                raise FileNotFoundException(f"File not found: {file_path}")

            file_extension = path.suffix.lower()
            if file_extension not in self.supported_formats:
                raise UnsupportedFormatException(f"Unsupported file format: {file_extension}")

            logger.info(f"Processing {path.name}")
            return self.supported_formats[file_extension](file_path)

        except (FileNotFoundException, UnsupportedFormatException):
            raise
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            raise DocumentProcessingException(f"Failed to process file {file_path}: {str(e)}")

    def read_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
            logger.debug(f"Extracted {len(text)} characters from DOCX")
            return text
        except Exception as e:
            logger.error(f"DOCX reading error: {str(e)}")
            raise

    def read_pdf(self, file_path: str) -> str:
        """Extract text from PDF file including tables and images."""
        try:
            # Extract regular text
            reader = PdfReader(file_path)
            text = ''
            for i, page in enumerate(reader.pages):
                text += page.extract_text()
                logger.debug(f"Processed PDF page {i+1}/{len(reader.pages)}")

            # Extract tables
            try:
                tables_text = self.extract_tables_from_pdf(file_path)
                if tables_text and tables_text != "No tables found":
                    text += f"\n\n[EXTRACTED TABLES]\n{tables_text}"
                    logger.info("Successfully extracted tables from PDF")
            except Exception as e:
                logger.warning(f"Table extraction failed: {str(e)}")

            # Extract text from images using OCR
            try:
                ocr_text = self.extract_text_from_images_in_pdf(file_path)
                if ocr_text.strip():
                    text += f"\n\n[OCR TEXT FROM IMAGES]\n{ocr_text}"
                    logger.info("Successfully extracted text from images in PDF")
            except Exception as e:
                logger.warning(f"OCR extraction failed: {str(e)}")

            return text
        except Exception as e:
            logger.error(f"PDF reading error: {str(e)}")
            raise

    def read_txt(self, file_path: str) -> str:
        """Read text file with proper encoding handling."""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Could not read file with any of the encodings: {encodings}")

    def read_excel(self, file_path: str) -> str:
        """Convert Excel file to string representation."""
        try:
            df = pd.read_excel(file_path)
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"Excel reading error: {str(e)}")
            raise

    def read_csv(self, file_path: str) -> str:
        """Convert CSV file to string representation."""
        try:
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        except Exception as e:
            logger.error(f"CSV reading error: {str(e)}")
            raise

    def read_ppt(self, file_path: str) -> str:
        """Extract text from PowerPoint file."""
        try:
            prs = Presentation(file_path)
            texts = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                texts.append(f"[Slide {i+1}]\n" + '\n'.join(slide_text))
                logger.debug(f"Processed slide {i+1}/{len(prs.slides)}")
            return '\n\n'.join(texts)
        except Exception as e:
            logger.error(f"PowerPoint reading error: {str(e)}")
            raise

    def read_md(self, file_path: str) -> str:
        """Read Markdown file."""
        return self.read_txt(file_path)

    def read_sqlite(self, file_path: str) -> str:
        """Extract data from SQLite database."""
        try:
            engine = create_engine(f'sqlite:///{file_path}')
            with engine.connect() as conn:
                tables = pd.read_sql_query(
                    "SELECT name FROM sqlite_master WHERE type='table'", conn)
                result = {}
                for table in tables['name']:
                    df = pd.read_sql_table(table, conn)
                    result[table] = df.to_string(index=False)
                    logger.debug(f"Processed table: {table}")
                return str(result)
        except Exception as e:
            logger.error(f"SQLite reading error: {str(e)}")
            raise

    def extract_tables_from_pdf(self, file_path: str) -> str:
        """Extract tables from PDF using Camelot and Tabula."""
        try:
            tables_texts = []
            # Try Camelot first
            try:
                camelot_tables = camelot.read_pdf(file_path, pages='all', flavor='stream')
                for i, table in enumerate(camelot_tables):
                    tables_texts.append(f"[Camelot Table {i+1}]\n{table.df.to_string(index=False)}")
                logger.info(f"Extracted {len(camelot_tables)} tables using Camelot")
            except Exception as e:
                logger.warning(f"Camelot failed: {str(e)}")

            # Try Tabula as fallback
            try:
                tabula_tables = tabula.read_pdf(file_path, pages='all', multiple_tables=True)
                for i, df in enumerate(tabula_tables):
                    tables_texts.append(f"[Tabula Table {i+1}]\n{df.to_string(index=False)}")
                logger.info(f"Extracted {len(tabula_tables)} tables using Tabula")
            except Exception as e:
                logger.warning(f"Tabula failed: {str(e)}")

            return '\n\n'.join(tables_texts) if tables_texts else "No tables found"
        except Exception as e:
            logger.error(f"Error extracting tables from PDF: {str(e)}")
            raise

    def extract_images_from_pdf(self, file_path: str) -> List[np.ndarray]:
        """Extract images from PDF pages as numpy arrays."""
        try:
            images = []
            pages = convert_from_path(file_path)
            for i, page in enumerate(pages):
                img = np.array(page)
                images.append(img)
                logger.debug(f"Extracted image from page {i+1}")
            return images
        except Exception as e:
            logger.error(f"Error extracting images from PDF: {str(e)}")
            raise

    def ocr_image(self, image: np.ndarray) -> str:
        """Perform OCR on a numpy image array."""
        try:
            pil_img = Image.fromarray(image)
            text = pytesseract.image_to_string(pil_img)
            logger.debug(f"OCR extracted {len(text)} characters")
            return text
        except Exception as e:
            logger.error(f"OCR error: {str(e)}")
            raise

    def extract_text_from_images_in_pdf(self, file_path: str) -> str:
        """Extract text from images embedded in PDF using OCR."""
        try:
            images = self.extract_images_from_pdf(file_path)
            ocr_texts = [self.ocr_image(img) for img in images]
            return '\n\n'.join(ocr_texts)
        except Exception as e:
            logger.error(f"Error extracting text from images in PDF: {str(e)}")
            raise
